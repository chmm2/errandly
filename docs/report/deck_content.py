"""Slide content for the Review-2 deck."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_deck import (
    ACCENT_FILL, BODY_FONT, GREY, INK, M, OUT, RULE, TITLE_FONT, W,
    bullets, figure, panel, slide, table, textbox,
)

TITLE = ("Errandly: A Trust-Aware, Fraud-Resistant Platform for "
         "Campus Errands, Micro-Delivery and Commerce")

LIT_A = [
    ["[1] Jiang et al.\nACM Comput. Surv.\n2016",
     "Formalise how trust is computed and propagated across a social graph, not "
     "just between two directly connected users.",
     "Comparative survey of graph-simplification and graph-analogy trust "
     "propagation.",
     "+ Makes hop-distance trust decay a tunable design problem\n"
     "− Survey only; no matching, fraud or campus application",
     "Gap: decay depends on distance alone.\n"
     "Ours: trust also decays with the AGE of the newest edge, so a friendship "
     "made yesterday buys almost nothing today."],
    ["[2] Chiou & Tu\nIEEE Access\n2020",
     "Let ride-hailing riders weight a driver's ratings by how close the rater "
     "is to them.",
     "Social-network scoring with cryptographic rater privacy; built and tested "
     "on Android.",
     "+ First to fold the user's own social graph into trust scoring\n"
     "− Rating weighting only; no price or fraud signal",
     "Gap: closeness raises a rater's weight.\n"
     "Ours: closeness LOWERS it once a runner's ratings concentrate inside their "
     "own circle — the farming case."],
    ["[3] Fu & Liu\nComputer Journal\n2021",
     "Assign crowdsourcing tasks to the most trusted worker, not the nearest or "
     "cheapest.",
     "T-Aware: trust scored from platform history, optimised jointly with task "
     "cost.",
     "+ Treats trust as a formal objective in assignment\n"
     "− Trust from platform history only; no peer graph; not campus-scoped",
     "Gap: no notion of who actually knows whom.\n"
     "Ours: a real friendship graph supplies trust, and every ranking term is "
     "expressed in one unit — metres."],
    ["[4] Cheng, Chen & Ye\nIEEE ICDE\n2019",
     "Co-assign workers who have collaborated successfully before.",
     "Cooperation-aware assignment; cooperativeness scored from interaction "
     "history.",
     "+ Measurable quality gains from using social history\n"
     "− The policy decides who transacts with whom — never examined",
     "Gap: the policy shapes the very history it later consumes.\n"
     "Ours: this feedback loop is our central research problem (G5)."],
    ["[5] Escrow-based P2P\npayment system\n2026",
     "Reduce fraud in peer-to-peer payments using a trusted third party.",
     "Object-oriented design; third-party escrow; PHP/MySQL implementation.",
     "+ Structured, auditable release of funds\n"
     "− Generic e-commerce; balances stored and mutated in place",
     "Gap: a stored balance can drift from its own history.\n"
     "Ours: append-only ledger, balances derived, and a uniqueness rule so a "
     "redelivered settlement collides instead of paying twice."],
]

LIT_B = [
    ["[6] Sarpal et al.\narXiv:2310.04367\n2023",
     "Detect mispriced or manipulated listings at scale without manual review.",
     "MoatPlus: ensemble of unsupervised models over historical and nearby "
     "prices sets a price bound.",
     "+ Price consensus needs no receipts; +46.6% precision in risky categories\n"
     "− Built for catalogue SKUs, not peer-reported cash spend",
     "Gap: one price per item ignores that the same item costs different amounts "
     "at different shops.\n"
     "Ours: a PER-STORE reference, shrunk toward the campus median by how much "
     "evidence supports it."],
    ["[7] Maranzato & Pereira\nIEEE LA-Web\n2009",
     "Detect sellers gaming reputation systems in electronic marketplaces.",
     "Behavioural and transaction features — pricing, feedback patterns — from "
     "real marketplace data.",
     "+ Transaction patterns flag reputation gaming without manual review\n"
     "− Detects fraud after the fact; never fed back into matching",
     "Gap: detection and matching are disconnected.\n"
     "Ours: an open flag is a bounded ranking penalty — suspicion demotes a "
     "runner at once, but can never exclude them."],
    ["[8] Kamvar et al.\nWWW\n2003",
     "Compute a global trust value for every peer in a P2P network.",
     "EigenTrust: transitive propagation of local trust, power iteration to a "
     "global vector.",
     "+ Canonical and fully distributed\n"
     "− Authors themselves note vulnerability to collusive groups rating each "
     "other up",
     "Gap: the collusive-group weakness is acknowledged, not solved.\n"
     "Ours: rings found as directed payment cycles among mutual friends "
     "(Tarjan's algorithm)."],
    ["[9] Viswanath et al.\nACM SIGCOMM\n2010",
     "Unify and evaluate social-graph defences against fake identities.",
     "Comparative analysis of SybilGuard, SybilLimit and SybilRank over real "
     "topologies.",
     "+ Shows these schemes are really community detection\n"
     "− Effectiveness rests entirely on the trust assumption holding",
     "Gap: the graph is treated as OBSERVED, never as produced by the platform.\n"
     "Ours: our graph is produced by our own routing, so its evidence must be "
     "conditioned on the policy that created it."],
    ["[10] Perdomo et al.\nICML\n2020",
     "Formalise settings where deploying a model changes the distribution it "
     "predicts.",
     "Performative prediction: risk minimisation over a model-dependent "
     "distribution.",
     "+ Exactly the right apparatus for a self-influencing system\n"
     "− Developed for prediction and pricing; never applied to fraud evidence",
     "Gap: nobody applies it to a fraud graph the platform itself generated.\n"
     "Ours: the routing policy becomes the null hypothesis against which "
     "in-group activity is judged."],
]

REFS = [
    "[1]   W. Jiang, G. Wang, Md. Z. A. Bhuiyan and J. Wu, “Understanding graph-based trust evaluation in online social networks,” ACM Computing Surveys, vol. 49, no. 1, 2016.",
    "[2]   S.-Y. Chiou and T.-Y. Tu, “A trusted mobile ride-hailing evaluation system with privacy and authentication,” IEEE Access, vol. 8, pp. 61929–61942, 2020.",
    "[3]   F. Donglai and L. Yanhua, “Trust-aware task allocation in collaborative crowdsourcing model,” The Computer Journal, vol. 64, no. 6, 2021.",
    "[4]   P. Cheng, L. Chen and J. Ye, “Cooperation-aware task assignment in spatial crowdsourcing,” in Proc. IEEE ICDE, 2019, pp. 1442–1453.",
    "[5]   “An escrow-based peer-to-peer online payment system for fraud reduction,” Methods in Science and Technology Studies, 2026.",
    "[6]   A. Sarpal, Q. Kang, F. Huang, Y. Song and L. Wan, “A marketplace price anomaly detection system at scale,” arXiv:2310.04367, 2023.",
    "[7]   R. P. Maranzato and A. M. Pereira, “Feature extraction for fraud detection in electronic marketplaces,” in Proc. IEEE LA-Web, 2009.",
    "[8]   S. D. Kamvar, M. T. Schlosser and H. Garcia-Molina, “The EigenTrust algorithm for reputation management in P2P networks,” in Proc. WWW, 2003, pp. 640–651.",
    "[9]   B. Viswanath, A. Post, K. P. Gummadi and A. Mislove, “An analysis of social network-based Sybil defenses,” in Proc. ACM SIGCOMM, 2010, pp. 363–374.",
    "[10] J. Perdomo, T. Zrnic, C. Mendler-Dünner and M. Hardt, “Performative prediction,” in Proc. ICML, 2020, pp. 7599–7609.",
    "[11] R. E. Tarjan, “Depth-first search and linear graph algorithms,” SIAM J. Computing, vol. 1, no. 2, pp. 146–160, 1972.",
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 0

    # ── 1 title ──────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    textbox(s, M, 1.35, W, 0.3, "B.Tech. Computer Science and Engineering   ·   "
            "BCSE497J  Project-I   ·   Review 2", 11, bold=True, color=GREY)
    textbox(s, M, 1.85, W, 0.7, "ERRANDLY", 46, bold=True, color=INK, font=TITLE_FONT)
    textbox(s, M, 2.70, 10.6, 1.0,
            "A Trust-Aware, Fraud-Resistant Platform for Campus Errands, "
            "Micro-Delivery and Commerce", 19, color=INK, font=TITLE_FONT, spacing=1.15)
    textbox(s, M, 4.30, 5.6, 1.4,
            "Team\nChris Martin Mattam   (23BCE0743)\n"
            "Sanskriti Sajlal   (23BCE0832)\nUjjwal Gogoi   (23BDS0335)",
            12.5, color=INK, spacing=1.3)
    textbox(s, 6.6, 4.30, 6.0, 1.4,
            "Faculty guide\nDr. Ranjithkumar S\n"
            "Assistant Professor, SCOPE\nProject ID: 20226UG03",
            12.5, color=INK, spacing=1.3)
    textbox(s, M, 6.55, W, 0.3,
            "School of Computer Science and Engineering  ·  Vellore Institute of "
            "Technology, Vellore  ·  September 2026", 10, color=GREY)

    # ── 2 approval ───────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Review 2", "Approval Mail From Guide",
              "Approval to be obtained by e-mail before the review on the points "
              "below.", page=n + 1)
    bullets(s, M, 2.20, 6.0, 2.0, [
        "Aim, objectives and problem statement",
        "Research gap established from the literature survey",
        "Proposed architecture and module breakdown",
        "Target outcome: Scopus-indexed conference paper",
    ], size=13)
    panel(s, 7.0, 2.10, 5.7, 3.4, "PASTE APPROVAL MAIL SCREENSHOT HERE", "")

    # ── 3 aim ────────────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 2.3", "Aim", page=n + 1)
    textbox(s, M, 1.45, W, 1.3,
            "To design, develop and validate Errandly — a campus-verified platform "
            "on which VIT students earn through paid peer-run errands and resell "
            "pre-owned essentials inside a closed, trust-verified network, with every "
            "transaction settled through an escrow wallet — so that campus errands and "
            "resale happen safely, without informal cash deals or outside gig-platform "
            "markups.", 15, color=INK, spacing=1.25)
    for i, (h, t) in enumerate([
        ("Verified, not open",
         "Only campus-email-verified VIT students can post or accept."),
        ("Escrow, not trust-me",
         "Payment is held by the app and released only on confirmed delivery."),
        ("Commerce, not just delivery",
         "A peer resale marketplace runs on the same rails as errands."),
        ("Trusted, and still checkable",
         "Routing uses who you know — without blinding our own fraud detection."),
    ]):
        panel(s, M + i * 3.06, 3.10, 2.88, 1.35, h, t, hsize=12.5, tsize=10.5)
    textbox(s, M, 4.75, W, 0.8,
            "Routing on who you know, and detecting collusion among the people you "
            "know, pull against each other: the harder we favour friends, the more "
            "in-group activity appears for innocent reasons. Resolving that tension is "
            "the research contribution.", 12.5, italic=True, color=GREY, spacing=1.2)

    # ── 4 objectives ─────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 2.3", "Objectives", page=n + 1)
    bullets(s, M, 1.50, W, 5.0, [
        "O1  Errand & marketplace core | four errand types — shopping, food, parcel, "
        "main-gate — plus a peer resale marketplace on one FastAPI modular monolith.",
        "O2  Identity & verification | campus-email signup with OTP; reputation built "
        "only from completed transactions.",
        "O3  Trust-first matching | rank nearby Redis-GEO runners by social proximity, "
        "reputation and integrity standing; ties broken by distance; race-guarded accept.",
        "O4  Escrow wallet | hold the requester's payment at posting, release on "
        "confirmed delivery, settle unpriced items against price consensus.",
        "O5  Auditable lifecycle | event-source every state transition through a "
        "transactional outbox so any errand or payment is reconstructable.",
        "O6  Multi-signal fraud detection | price inflation per claim, rating farming "
        "per runner, collusion rings per group — fed back as a bounded penalty, not a ban.",
        "O7  Policy-conditioned evidence | log the full candidate set and every ranking "
        "term, and use the routing policy itself as the null hypothesis for collusion.",
        "O8  Validation | pilot with a VIT cohort measuring completion, dispute and "
        "time-to-match, plus a simulation of the detector's false-positive behaviour.",
    ], size=12, gap=8)
    textbox(s, M, 6.45, W, 0.4,
            "SDG 8 Decent Work and Economic Growth (primary)  ·  SDG 11 Sustainable "
            "Cities  ·  SDG 12 Responsible Consumption      |      Outcome: working "
            "prototype + Scopus-indexed conference paper", 10.5, color=GREY)

    # ── 5 abstract ───────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Overview", "Abstract", page=n + 1)
    bullets(s, M, 1.50, W, 5.2, [
        "Errandly is a campus-verified platform where VIT students post everyday "
        "errands — shopping lists, food orders, parcel pickups, main-gate collections — "
        "and resell pre-owned hostel items to each other.",
        "Every transaction settles through an escrow wallet: the requester's money is "
        "held when the task is posted and released only once delivery is confirmed, so "
        "neither side has to trust the other first.",
        "Errands go to nearby runners by a ranking that combines walking distance, "
        "social closeness from a friendship graph, and reputation earned from completed "
        "transactions.",
        "Three detectors guard the platform: one checks each reported price against "
        "what that item costs at that shop; one checks whether a runner's ratings come "
        "only from their own friends; one looks for a collusion ring.",
        "A collusion ring is a group of friends who keep paying each other in a closed "
        "loop, so the same money circles the same few people while every lap "
        "manufactures reputation and errand history at no real cost.",
        "The problem we found: because the matcher offers errands to friends first, "
        "friends transact far more than chance predicts — which is exactly what the "
        "ring detector treats as evidence. The platform manufactures its own suspicion.",
        "Our fix: work out how much in-group activity the matcher itself predicted, and "
        "count only the unexplained excess as evidence. Honest friend groups clear; a "
        "real ring still shows.",
    ], size=12, gap=9)

    # ── 6 problem statement ──────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 2.4", "Problem Statement",
              "The real-world gap, and the engineering problem inside it.", page=n + 1)
    for i, (big, small) in enumerate([
        ("21–24%", "average cut food-delivery apps take from every order"),
        ("₹17.58", "flat platform fee stacked on top, per order"),
        ("0", "campus platforms verify both sides before money moves"),
        ("Every sem.", "usable furniture, cycles and books discarded by leavers"),
    ]):
        panel(s, M + i * 3.06, 2.15, 2.88, 1.30, big, small, hsize=20, tsize=10.5)
    bullets(s, M, 3.75, W, 2.6, [
        "Engineering problem | build a closed, campus-verified platform that matches "
        "errands and resale listings to nearby trusted peers, holds funds neutrally "
        "until delivery is confirmed, and catches price and reputation abuse where no "
        "receipts exist.",
        "Research problem | routing on social trust and detecting collusion among the "
        "socially connected pull against each other. The stronger the social preference, "
        "the more in-group concentration appears for innocent reasons, and the less of "
        "it survives as evidence.",
        "The task | not to detect collusion, but to separate the platform's own "
        "contribution to a pattern from the participants' — and to establish how strong "
        "a trust boost can be before that separation becomes impossible.",
    ], size=12, gap=9)

    # ── 7-8 literature ───────────────────────────────────────────────────────
    for part, rows in (("1 of 2", LIT_A), ("2 of 2", LIT_B)):
        n += 1
        s = slide(prs, f"Section 2.1  ·  {part}", "Literature Review",
                  "Ten papers. A paper appears only where we implement something that "
                  "answers it. The last column is the gap it leaves and what we do "
                  "about it.", page=n + 1)
        table(s, M, 2.10, W,
              ["Paper", "Objective", "Methodology", "Pros / Cons",
               "Gap left  →  our response"],
              rows, widths=[1.75, 2.30, 2.30, 2.70, 3.04],
              hdr_size=9.1, body_size=8.2, row_h=0.80)

    # ── 9 research gap ───────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 2.2", "Research Gap",
              "Four gaps shape the platform. The fifth emerged during implementation "
              "and is the research contribution.", page=n + 1)
    table(s, M, 2.15, W, ["#", "Gap", "What it means"], [
        ["G1", "Trust is assumed, not engineered",
         "Campus apps rely on goodwill or a star rating. Nothing holds payment until a "
         "task is actually done, so all the risk sits with whoever moves first."],
        ["G2", "Commerce and tasks live apart",
         "Marketplace apps solve resale; delivery apps solve tasks. No platform lets one "
         "verified campus identity do both under one wallet and one reputation."],
        ["G3", "Price integrity without receipts",
         "Campus shops issue no receipts. Existing price-anomaly work [6] assumes one "
         "price per item — but the same puff costs ₹23 at one canteen and ₹30 at another."],
        ["G4", "Detection never reaches the matcher",
         "Marketplace fraud detection [7] is retrospective. A runner under active "
         "suspicion still receives exactly as many offers as before."],
        ["G5", "The platform manufactures its own fraud evidence",
         "Trust-aware assignment [3][4] promotes friends. Collusion detection [8][9] "
         "reads friends transacting as evidence of a ring. The first CAUSES the second."],
    ], widths=[0.55, 3.30, 8.24], hdr_size=9.5, body_size=9.5, row_h=0.86)
    textbox(s, M, 6.55, W, 0.5,
            "Prior work assumes an EXOGENOUS graph — something the platform observes. "
            "Ours is ENDOGENOUS, generated by our own routing decisions. That is the "
            "gap, and performative prediction [10] is the apparatus nobody has applied "
            "to it.", 11.5, color=GREY, spacing=1.15)

    # ── 10 what a ring is ────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Concept", "What a Collusion Ring Actually Is",
              "The attack our graph layer exists to catch — and why structure alone "
              "cannot catch it.", page=n + 1)
    figure(s, "fig_social.png", 0.62, 2.05, 7.4)
    panel(s, 8.25, 2.05, 4.45, 1.55, "The attack",
          "Three friends take turns running errands for each other and paying each "
          "other. No goods really move. The same money goes round the circle.",
          hsize=13, tsize=11)
    panel(s, 8.25, 3.75, 4.45, 1.45, "Why they bother",
          "Each lap manufactures a completed errand and a five-star rating — "
          "reputation and history bought at no real cost.", hsize=13, tsize=11)
    panel(s, 8.25, 5.32, 4.45, 1.45, "Why it is hard",
          "An honest friend group and a ring look IDENTICAL structurally — both fully "
          "closed. The discriminator is economic: does value circulate and return?",
          hsize=13, tsize=11)

    # ── 11 what we add ───────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Contribution", "What We Are Adding",
              "The problem, the fix and the limit — in plain terms.", page=n + 1)
    panel(s, M, 2.05, 3.88, 2.05, "1.  The problem we found",
          "Our matcher offers errands to friends FIRST.\n"
          "So friends transact far more than chance.\n"
          "So the ring detector sees 'suspicious' concentration.\n"
          "The platform manufactures its own evidence.", hsize=13, tsize=11)
    panel(s, M + 4.10, 2.05, 3.88, 2.05, "2.  The fix",
          "Ask what the matcher itself PREDICTED would happen.\n"
          "Count only the unexplained excess as evidence.\n"
          "Honest friends follow the policy → excess ≈ 0.\n"
          "A ring pairs up beyond it → excess stays large.", hsize=13, tsize=11)
    panel(s, M + 8.20, 2.05, 3.88, 2.05, "3.  The limit we found",
          "Boost friends hard enough and the policy expects\n"
          "~everything to stay in-group. Nothing can then be\n"
          "surprising, so collusion becomes undetectable —\n"
          "not because it hid, but because routing did its job.",
          hsize=13, tsize=11)
    panel(s, M, 4.40, W, 1.05, "4.  The remedy",
          "Offer a small fraction of errands (ε ≈ 5%) with friendship ignored entirely "
          "— distance only, no hop ceiling. Those rounds are the control group: the "
          "only errands whose outcome was not already decided by who is friends with "
          "whom. It is a configuration value, not a redesign.", hsize=13, tsize=11.5)
    textbox(s, M, 5.75, W, 0.9,
            "Why it matters beyond us:  any platform that routes on trust and then "
            "mines the resulting interactions for fraud is reading its own routing "
            "decisions back as evidence. We believe this has not been stated before.",
            12, italic=True, color=GREY, spacing=1.2)

    # ── 12 project plan ──────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 2.5", "Project Plan",
              "Six sprints overlapping by about a week, July to October 2026. Ordering "
              "is by dependency: settlement records are the payment edges the graph "
              "projects, and ring detection runs over those edges.", page=n + 1)
    figure(s, "fig_gantt.png", 1.30, 2.20, 10.7)

    # ── 13 architecture ──────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 4.1", "Architecture",
              "A modular monolith over polyglot persistence. Bold outlines mark the "
              "modules carrying the contribution; Neo4j is a derived read model, "
              "rebuildable from PostgreSQL at any time.", page=n + 1)
    figure(s, "fig_architecture.png", 3.15, 2.05, 7.0)

    # ── 14 modules ───────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 4.2", "Modules",
              "Six domain modules with enforced boundaries — a module reaches another "
              "only through its service interface, never its tables.", page=n + 1)
    table(s, M, 2.15, W, ["Module", "Responsibility", "Store"], [
        ["Auth & Identity", "Registration, OTP verification, sessions, recovery, role "
         "and campus tenancy.", "PostgreSQL"],
        ["Errands & Matching", "Creation, candidate generation, ranking, tiered "
         "offering, contention control, offer log.", "PostgreSQL + Redis"],
        ["Ledger & Escrow", "Holds, settlement, refunds, clawbacks. Append-only; "
         "balances derived, never stored.", "PostgreSQL"],
        ["Social Graph", "Friendship and payment edges projected from Postgres; trust "
         "as maturity-weighted path length.", "Neo4j (derived)"],
        ["Fraud Detection", "Price inflation, rating farming, collusion rings, LLM "
         "corroboration, admin review.", "PostgreSQL + Neo4j"],
        ["Chat & Notify", "Per-errand conversation, live status streaming, push "
         "notification.", "MongoDB + Redis"],
    ], widths=[2.20, 7.49, 2.40], hdr_size=10, body_size=10, row_h=0.62)

    # ── 15-19 module figures ─────────────────────────────────────────────────
    for eyebrow, title, lede, fig, x, w in [
        ("Section 4.2.1", "Trust-Aware Matching",
         "Every ranking term is in metres, so the terms stay comparable. The exploring "
         "branch is the control group: the only rounds not already shaped by friendship.",
         "fig_matching.png", 3.60, 6.1),
        ("Section 4.2.3", "Escrow and the Append-Only Ledger",
         "The balance is derived — credits less debits — never stored, never mutated. A "
         "redelivered settlement collides on a uniqueness constraint instead of paying "
         "twice.", "fig_escrow.png", 2.95, 7.4),
        ("Section 4.2.4", "Social Graph and Collusion Rings",
         "Structure alone cannot separate a genuine friend group from a ring; both are "
         "fully closed. The discriminator is economic — does value circulate and return "
         "to its source?", "fig_social.png", 2.75, 7.8),
        ("Section 4.2.5", "Multi-Signal Fraud Detection",
         "Cheap deterministic arithmetic decides who to examine; the language model is "
         "consulted only afterwards and its verdict is advisory. No detector punishes "
         "on its own authority.", "fig_fraud.png", 2.90, 7.5),
        ("Section 4.2.6", "End-to-End Sequence",
         "The escrow hold is synchronous, because a requester must be told at once if "
         "their balance is short. Settlement is asynchronous; dashed arrows mark that "
         "boundary.", "fig_sequence.png", 3.05, 7.2),
    ]:
        n += 1
        s = slide(prs, eyebrow, title, lede, page=n + 1)
        figure(s, fig, x, 2.15, w)

    # ── 20 results ───────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Evaluation", "Experiments and Results",
              "Simulation over a synthetic campus: 60 students, 6 000 errands, six "
              "honest friend groups sharing hostel blocks, one genuine three-person "
              "ring.", page=n + 1)
    table(s, M, 2.10, W, ["Experiment", "Setup", "Result"], [
        ["Naive vs. policy-conditioned",
         "At the deployed trust setting (1500 m boost)",
         "Naive in-group concentration flags 6 of 6 HONEST groups. Policy-conditioned "
         "excess flags 0 of 6 — and still identifies the ring."],
        ["Detectability limit",
         "Trust boost swept 0 → 2500 m, detector fixed",
         "As the boost rises the policy expects ~99% in-group routing. Excess falls to "
         "z = 2.3 and the ring becomes undetectable in principle."],
        ["Deliberate exploration",
         "Socially-blind fraction ε varied at the deployed setting",
         "ε = 0%: ring missed (z = 2.3).  ε = 5%: caught (z = 3.8).  ε = 10%: "
         "z = 4.5. Honest groups stay clear throughout."],
        ["Co-located ring (stress test)",
         "Ring members placed at the same desk",
         "Still caught at 85% loyalty (z = 3.7). A ring hides only by RESTRAINT — which "
         "rate-limits how much it can farm."],
        ["Implementation",
         "Full stack under Docker Compose, automated test suite",
         "Escrow, tiered offering, three detectors, admin review console, offer log and "
         "ε-exploration running end to end; 214 tests passing."],
    ], widths=[2.45, 3.30, 6.34], hdr_size=9.5, body_size=9.2, row_h=0.80)

    # ── 21 conclusion ────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 5", "Conclusion", page=n + 1)
    bullets(s, M, 1.55, W, 3.6, [
        "Errandly is a working campus platform: escrow-backed settlement, trust-aware "
        "matching, a three-signal fraud pipeline and an admin review path, on mobile "
        "and web.",
        "The finding that matters is structural. A trust-aware matcher and a "
        "graph-based collusion detector — each recommended independently by the "
        "literature — undermine one another when deployed together.",
        "Making the routing policy the null hypothesis corrects this, and replaces "
        "hand-tuned thresholds with an actual test statistic.",
        "A strong enough trust boost makes collusion undetectable in principle. "
        "Deliberate exploration buys that observability back, and it is a configuration "
        "value rather than a redesign.",
    ], size=12.5, gap=10)
    panel(s, M, 5.05, W, 1.55, "Project-II",
          "Peer resale marketplace implementation  ·  field study with a live VIT "
          "cohort  ·  real-money settlement through a licensed gateway  ·  submission "
          "of the policy-conditioned detection result to a Scopus-indexed venue",
          hsize=13, tsize=11.5)

    # ── 22 references ────────────────────────────────────────────────────────
    n += 1
    s = slide(prs, "Section 5", "References", page=n + 1)
    tb = s.shapes.add_textbox(Inches(M), Inches(1.55), Inches(W), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ref in enumerate(REFS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6); p.line_spacing = 1.05
        r = p.add_run(); r.text = ref
        r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.color.rgb = INK

    prs.save(OUT)
    print(f"written: {OUT}")
    print(f"  slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")


if __name__ == "__main__":
    build()
