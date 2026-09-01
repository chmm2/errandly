"""Sanskriti's Review-2 slides.

    cd docs/report && python deck_sanskriti.py

Reuses the layout helpers from deck_chris so the two decks are visually
identical — same margins, same type, same plain white. Only the footer differs.

Every slide: the gap in prior work -> what we built -> what it changes.
"""

from __future__ import annotations

import pathlib

import deck_chris
from deck_chris import GREY, INK, M, W, bullets, figure, rule, table, text
from pptx import Presentation
from pptx.util import Inches, Pt

# Same helpers, her name in the footer.
deck_chris.FOOTER = "Errandly  ·  Sanskriti Sajlal  ·  23BCE0832"
from deck_chris import slide  # noqa: E402  — imported after the footer is set

OUT = pathlib.Path(__file__).parent / "Errandly-Review2-Sanskriti.pptx"


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 0

    # ══════════════════════════════════════════════════════════════ 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, M, 1.55, W, 0.3,
         "BCSE497J  PROJECT-I   ·   REVIEW 2   ·   INDIVIDUAL CONTRIBUTION",
         size=11, bold=True, color=GREY)
    text(s, M, 2.05, 11.6, 1.5, "Trust-Aware Social Matching",
         size=36, bold=True, font="Cambria", spacing=1.12)
    text(s, M, 3.05, 11.6, 0.6,
         "Why a new friendship is worth less than an old one",
         size=16, color=INK, font="Cambria")
    rule(s, 3.95)
    text(s, M, 4.20, 11.6, 0.9,
         "Errandly: A Trust-Aware, Fraud-Resistant Platform for Campus Errands, "
         "Micro-Delivery and Commerce", size=13, color=GREY, spacing=1.2)
    text(s, M, 5.25, 6.0, 1.0, "Sanskriti Sajlal\n23BCE0832",
         size=13.5, spacing=1.35)
    text(s, 6.9, 5.25, 5.8, 1.0,
         "Guide: Dr. Ranjithkumar S\nAssistant Professor, SCOPE",
         size=13.5, spacing=1.35)
    text(s, M, 6.85, W, 0.3,
         "Vellore Institute of Technology, Vellore   ·   September 2026",
         size=10, color=GREY)

    # ══════════════════════════════════════════════════════════ 2 the scope
    n += 1
    s = slide(prs, "Scope", "What This Half Of The Project Covers",
              "Five mechanisms. Four are applied work. One is the contribution.",
              page=n + 1)
    table(s, M, 2.10, W, ["#", "Mechanism", "The gap it answers"],
          [["1", "Friendship graph, projected and rebuildable",
            "Trust-aware allocation [3] scores trust from platform history "
            "only — there is no peer graph at all"],
           ["2", "Trust over paths, up to four hops",
            "[1] formalises how trust decays along a path. Applied here."],
           ["3", "**Maturity — trust decays with the AGE of the newest link",
            "**[1] decays trust with distance only. None of the three papers "
            "models time."],
           ["4", "Tiered offering — two hops for 45 seconds",
            "[2] folds the social graph into rating weight, but it never "
            "reaches the assignment decision"],
           ["5", "Closure penalty on a sealed neighbourhood",
            "Structure alone cannot separate a friend group from a ring"]],
          widths=[0.4, 4.3, 7.43], hdr=11, body=11, rh=0.62)
    text(s, M, 6.25, W, 0.6,
         "Rows 1, 2, 4 and 5 apply what the literature already establishes. "
         "Row 3 is what is new.", size=12.5, bold=True)

    # ══════════════════════════════════════════════════════ 3 the problem
    n += 1
    s = slide(prs, "The setting", "Why Trust Belongs In Matching", page=n + 1)
    bullets(s, M, 2.05, W, 3.0, [
        "A market between strangers | money moves before the service is "
        "rendered, and the quality of that service is observed only by the two "
        "people involved.",
        "Distance is a poor way to choose | the nearest available student is "
        "not necessarily the one most likely to actually deliver.",
        "But trust on a campus is not uniform | students already know some of "
        "their neighbours, and an errand run by an acquaintance carries "
        "materially less risk than one run by a stranger.",
        "That structure is free information | it already exists, and the "
        "platform did not have to build it.",
    ], size=13.5, gap=13)
    rule(s, 5.35)
    text(s, M, 5.60, W, 0.9,
         "This layer has to do three things with it: represent it, turn it into "
         "a number, and stop it being gamed.\n"
         "Everything downstream — price checking, collusion detection, the "
         "escrow decisions — runs on errands this layer routed.",
         size=12.5, spacing=1.3)

    # ═══════════════════════════════════════════════════ 4 the three papers
    n += 1
    s = slide(prs, "Section 2.1", "What The Previous Work Does",
              "Three papers define the state of the art for this problem.",
              page=n + 1)
    table(s, M, 2.05, W,
          ["Paper", "What it establishes", "What it leaves open"],
          [["[1] Jiang et al.\nACM Computing\nSurveys, 2016",
            "Trust must DECAY along a path — a friend-of-a-friend-of-a-friend "
            "cannot count as much as a direct friend. The rate of decay is a "
            "design parameter.",
            "Decay is a function of DISTANCE ALONE. An edge is an edge: a "
            "friendship from three years ago and one from this morning "
            "contribute identically."],
           ["[2] Chiou & Tu\nIEEE Access, 2020",
            "A rating from someone close to you is worth more than one from a "
            "stranger. Implemented with cryptographic rater privacy.",
            "Stops at RATING WEIGHT. The social graph never reaches the "
            "assignment decision itself."],
           ["[3] Fu & Liu\nThe Computer\nJournal, 2021",
            "Trust as a FORMAL OBJECTIVE in task allocation, not an "
            "afterthought. Beats cost-optimal assignment on reliability.",
            "Trust comes only from PLATFORM HISTORY. No peer graph, so a new "
            "user is indistinguishable from any other new user."]],
          widths=[2.2, 4.9, 5.03], hdr=10.5, body=10.5, rh=1.05)
    rule(s, 5.95)
    text(s, M, 6.20, W, 0.8,
         "All three agree trust should inform assignment. None of them models "
         "TIME — a relationship either exists or it does not, and if it exists "
         "it counts in full from the instant it is created.",
         size=13, bold=True, spacing=1.25)

    # ══════════════════════════════════════════════════════════ 5 the gap
    n += 1
    s = slide(prs, "Section 2.2", "The Gap, And Why It Matters Here",
              "On a campus, creating a friendship is free and instant. A "
              "request, an accept, and an edge exists.", page=n + 1)
    text(s, M, 2.05, W, 0.6,
         "If that edge immediately confers full trust, the entire social layer "
         "can be acquired in an afternoon.", size=14, bold=True)

    text(s, M, 2.75, 5.7, 0.3, "TWO FORMS OF THE ATTACK", size=10.5, bold=True,
         color=GREY)
    bullets(s, M, 3.05, 5.7, 2.2, [
        "Targeted | find a student who posts many errands, befriend them, and "
        "collect a 1500 m head start on everything they post from that moment.",
        "Breadth | ratings are weighted by distinct rater, so repeating is "
        "worthless. The cheapest remaining evasion is to recruit eighty friends "
        "and have each rate once.",
    ], size=11.5, gap=10)

    text(s, 6.75, 2.75, 6.0, 0.3, "WHY THE SECOND MATTERS MOST", size=10.5,
         bold=True, color=GREY)
    text(s, 6.75, 3.05, 6.0, 2.2,
         "It is where the rest of the fraud work pushes an attacker.\n\n"
         "Every one of those eighty ties has to be created — and created "
         "recently.\n\n"
         "If new edges carry full weight, the defence has only moved the attack "
         "rather than stopped it.", size=12, spacing=1.35)

    rule(s, 5.55)
    text(s, M, 5.80, W, 0.6,
         "Stated as a gap: the literature decays trust across graph distance, "
         "and leaves it undecayed across time. On a platform where "
         "relationships are cheap to manufacture, that is the dimension an "
         "attacker moves in.", size=12.5, spacing=1.3)

    # ═══════════════════════════════════════════════════ 6 what we built
    n += 1
    s = slide(prs, "Section 4", "What We Built", page=n + 1)
    text(s, M, 1.95, W, 0.5,
         "trust   =   0.45 ^ (hops − 1)     ×     (1 − penalty)     ×     "
         "maturity", size=19, bold=True, font="Cambria")
    rule(s, 2.62)
    table(s, M, 2.85, W, ["Factor", "The question it answers", "Behaviour"],
          [["0.45 ^ (hops − 1)", "How far away are they?",
            "A direct friend scores 1.0. Each further hop keeps 45%."],
           ["1 − penalty", "Is this group sealed off from the rest of campus?",
            "Discounts trust flowing through a tightly closed neighbourhood."],
           ["**maturity", "**How new is the newest link on the path?",
            "**A brand-new friendship counts 0.40, maturing to 1.0 over "
            "thirty days."]],
          widths=[2.4, 4.6, 5.13], hdr=11, body=11, rh=0.56)
    text(s, M, 5.15, W, 0.3, "THE GRAPH ITSELF", size=10.5, bold=True,
         color=GREY)
    bullets(s, M, 5.45, W, 1.4, [
        "Friendships projected from PostgreSQL into Neo4j, each stamped with "
        "the date it was accepted",
        "A derived read model — nothing is authoritative there, and it can be "
        "rebuilt from PostgreSQL at any time",
        "Every read degrades to a neutral value if the graph is unavailable — "
        "matching becomes less targeted, never stops",
    ], size=11.5, gap=6)

    # ═════════════════════════════════════════════════ 7 matching depth
    n += 1
    s = slide(prs, "Mechanism 2", "How Far Trust Reaches", page=n + 1)
    figure(s, "fig_sk_hops.png", 1.65, 4.95)
    text(s, M, 6.72, W, 0.5,
         "Why stop at four: beyond four hops a path exists, but in practice it "
         "means nothing.", size=11.5, color=GREY)

    # ═══════════════════════════════════════ 8 the contribution — decay
    n += 1
    s = slide(prs, "Mechanism 3  ·  the contribution", "Friendship Decay",
              page=n + 1)
    figure(s, "fig_sk_decay.png", 1.55, 4.55)
    text(s, M, 6.28, W, 0.7,
         "The literature decays trust across graph distance. This decays it "
         "across edge age as well — so a relationship has to be established "
         "before it can be spent.", size=12.5, bold=True, spacing=1.25)

    # ══════════════════════════════════════ 9 three properties of the decay
    n += 1
    s = slide(prs, "Mechanism 3", "Three Deliberate Properties", page=n + 1)
    table(s, M, 2.00, W, ["Property", "Why"],
          [["It is the NEWEST edge on the path, not the average",
            "A chain of trust is only as established as its most recently "
            "created link. A two-year friendship reached through an edge made "
            "yesterday tells you about yesterday. Averaging would let one old "
            "friendship launder a brand-new one."],
           ["There is a floor of 0.40, not a zero",
            "Most new friendships are exactly what they appear to be. A new "
            "tie is discounted, never treated as worthless — the ordinary case "
            "must not be punished to inconvenience the rare one."],
           ["Old edges are NOT decayed",
            "Age makes a friendship stronger evidence, not weaker: a tie that "
            "has survived two years is more likely real than one from last "
            "week. Decaying it would penalise the whole ordinary population "
            "while costing an attacker nothing — rebuilding a stale network is "
            "free."]],
          widths=[3.6, 8.53], hdr=11, body=11, rh=1.02)
    text(s, M, 5.45, W, 0.5,
         "An unknown edge date counts in full: edges written before the date "
         "field existed are genuinely old, and treating missing data as "
         "suspicious would penalise the platform's earliest users for the "
         "platform's own gap.", size=11.5, color=GREY, spacing=1.25)
    rule(s, 6.15)
    text(s, M, 6.35, W, 0.4,
         "Six dedicated tests, including one that asserts recruiting a network "
         "today does not pay off today.", size=11.5, bold=True)

    # ═════════════════════════════════════════ 10 what the attack costs
    n += 1
    s = slide(prs, "Mechanism 3  ·  what it changes",
              "What A Same-Day Friendship Buys", page=n + 1)
    figure(s, "fig_sk_attack.png", 1.75, 3.55)
    table(s, M, 5.35, 8.4,
          ["", "Trust", "Head start"],
          [["Direct friend, established", "1.0000", "1500 m"],
           ["**Direct friend, made today", "**0.4200", "**630 m"],
           ["Two hops, newest link one day old", "0.1890", "284 m"]],
          widths=[4.4, 2.0, 2.0], hdr=11, body=11, rh=0.36)
    text(s, 9.15, 5.35, 3.6, 1.3,
         "Hop decay and maturity MULTIPLY.\n\n"
         "Freshness and distance compound rather than trading off.",
         size=11.5, color=GREY, spacing=1.3)

    # ═══════════════════════════════════════ 11 decisions and limitations
    n += 1
    s = slide(prs, "Design", "Decisions Worth Defending, And The Limits",
              page=n + 1)
    text(s, M, 2.00, 6.0, 0.3, "TWO DECISIONS", size=10.5, bold=True,
         color=GREY)
    bullets(s, M, 2.30, 6.0, 2.6, [
        "Closure, not the clustering coefficient | clustering divides by "
        "degree × (degree − 1), so a four-person ring scores about 0.5 while a "
        "ten-person hostel block scores 0.8. That ranks the threat exactly "
        "backwards. Closure is size-independent.",
        "Structure alone never discounts a direct friend | a close-knit friend "
        "group is the most ordinary thing on a campus, and the strongest honest "
        "trust signal we have.",
    ], size=11.5, gap=11)

    text(s, 7.05, 2.00, 5.7, 0.3, "LIMITATIONS", size=10.5, bold=True,
         color=GREY)
    bullets(s, 7.05, 2.30, 5.7, 3.4, [
        "The constants are unfitted — no live data supports 30 days, 0.40 or "
        "0.45 specifically",
        "Maturity delays an attack, it does not prevent one — a patient "
        "attacker who builds a network a month ahead pays nothing",
        "Trust is symmetric — friendship is an undirected edge, so A relying "
        "on B more than B relies on A cannot be expressed",
        "The graph can drift — it is rebuildable, but nothing detects "
        "divergence automatically",
    ], size=11.5, gap=8)
    rule(s, 5.90)
    text(s, M, 6.15, W, 0.4,
         "Fitting these constants against real campus data is Project-II work.",
         size=12, italic=True, color=GREY)

    # ══════════════════════════════════════════════════════ 12 summary
    n += 1
    s = slide(prs, "Summary", "Gaps Closed", page=n + 1)
    table(s, M, 2.05, W, ["Prior work", "The gap it leaves", "What we do"],
          [["[1] Jiang et al., ACM Computing Surveys, 2016",
            "Trust decays with graph distance. An edge's age is never modelled.",
            "**Trust also decays with the age of the newest link on the path"],
           ["[2] Chiou & Tu, IEEE Access, 2020",
            "The social graph informs rating weight, never the assignment "
            "decision",
            "Social proximity enters matching directly, in metres of effective "
            "distance"],
           ["[3] Fu & Liu, The Computer Journal, 2021",
            "Trust from platform history only — no notion of who knows whom",
            "A real friendship graph supplies trust, traversed up to four hops"]],
          widths=[3.1, 4.5, 4.53], hdr=11, body=11, rh=0.80)
    rule(s, 5.05)
    text(s, M, 5.30, W, 0.8,
         "The graph traversal and the hop decay are applied work.\n"
         "The contribution is the time dimension — a relationship must be "
         "established before it can be spent.", size=13.5, bold=True,
         spacing=1.3)
    text(s, M, 6.45, W, 0.4,
         "Both figures in this deck are generated from the constants the "
         "deployed system uses, so they cannot drift from the code.",
         size=11, italic=True, color=GREY)

    prs.save(OUT)
    print(f"written: {OUT}")
    print(f"  slides: {len(prs.slides._sldIdLst)}")


if __name__ == "__main__":
    build()
