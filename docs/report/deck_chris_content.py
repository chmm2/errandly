"""Slide content for Chris's Review-2 deck."""

from pptx import Presentation
from pptx.util import Inches, Pt

from deck_chris import (
    GREY, INK, M, OUT, W, bullets, figure, rule, slide, table, text,
)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 0

    # ══════════════════════════════════════════════════════════════ 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, M, 1.55, W, 0.3,
         "BCSE497J  PROJECT-I   ·   REVIEW 2   ·   INDIVIDUAL CONTRIBUTION",
         size=11, bold=True, color=GREY)
    text(s, M, 2.05, 11.6, 1.5,
         "Collusion Rings, and Evidence\nConditioned on Our Own Router",
         size=34, bold=True, font="Cambria", spacing=1.12)
    rule(s, 3.85)
    text(s, M, 4.10, 11.6, 0.9,
         "Errandly: A Trust-Aware, Fraud-Resistant Platform for Campus Errands, "
         "Micro-Delivery and Commerce", size=13, color=GREY, spacing=1.2)
    text(s, M, 5.20, 6.0, 1.0,
         "Chris Martin Mattam\n23BCE0743", size=13.5, spacing=1.35)
    text(s, 6.9, 5.20, 5.8, 1.0,
         "Guide: Dr. Ranjithkumar S\nAssistant Professor, SCOPE",
         size=13.5, spacing=1.35)
    text(s, M, 6.85, W, 0.3,
         "Vellore Institute of Technology, Vellore   ·   September 2026",
         size=10, color=GREY)

    # ══════════════════════════════════════════════════════ 2 what I own
    n += 1
    s = slide(prs, "Scope", "What This Half Of The Project Covers",
              "Five mechanisms. Each answers a specific gap in the papers we "
              "reviewed, and each ends in a consequence for the offender.",
              page=n + 1)
    table(s, M, 2.05, W,
          ["#", "Mechanism", "The gap it answers", "Consequence"],
          [["1", "Rating farming detection",
            "EigenTrust admits collusive groups defeat it, and leaves it unsolved",
            "Reputation discounted + 700 m"],
           ["2", "Collusion ring detection",
            "Social-graph defences assume the graph is observed, not produced",
            "1200 m + removed from each other's errands"],
           ["3", "The offer log",
            "Marketplace fraud work is retrospective; nothing records WHY a "
            "match was made",
            "None — it is evidence collection"],
           ["4", "Deliberate exploration",
            "Exploration is standard for unbiased evaluation, never for fraud "
            "observability",
            "None — it changes what data exists"],
           ["5", "LLM corroboration",
            "Structure and money answer who and how much, never what for",
            "None — advisory, attached to an existing flag"]],
          widths=[0.4, 2.5, 5.4, 3.83], hdr=10.5, body=10.5, rh=0.62)
    text(s, M, 6.35, W, 0.5,
         "Only the first two can penalise anyone. The other three exist so that "
         "the first two accuse the right people.",
         size=12, italic=True, color=GREY)

    # ══════════════════════════════════════════════════ 3 rating farming
    n += 1
    s = slide(prs, "Mechanism 1", "Rating Farming",
              "The attack: get your own friends to rate you, and buy a "
              "reputation you never earned.", page=n + 1)
    text(s, M, 2.02, 5.5, 0.3, "THE GAP", size=10.5, bold=True, color=GREY)
    bullets(s, M, 2.32, 5.55, 2.2, [
        "Kamvar et al., WWW 2003 | EigenTrust is the canonical reputation "
        "algorithm. Its own authors state it is defeated by groups who rate "
        "each other up — named, not solved.",
        "Chiou & Tu, IEEE Access 2020 | weights a rating by how close the rater "
        "is to you. Closeness only ever RAISES weight.",
    ], size=11.5, gap=10)
    text(s, M, 4.62, 5.55, 0.9,
         "A reputation built entirely inside your own circle is the signature of "
         "manipulation, not a mark of quality. Prior work treats it as the "
         "opposite.", size=11.5, italic=True, color=GREY)

    text(s, 6.55, 2.02, 6.2, 0.3, "WHAT WE BUILT", size=10.5, bold=True, color=GREY)
    bullets(s, 6.55, 2.32, 6.2, 2.6, [
        "One rater, one vote | ten ratings from one friend count as one opinion",
        "Friend ratings discounted | the more your ratings come from your own "
        "circle, the less each is worth",
        "A few ratings barely move you | you have to earn your way up",
    ], size=11.5, gap=9)
    text(s, 6.55, 4.72, 6.2, 0.3, "THEN THREE QUESTIONS — ANY YES RAISES A FLAG",
         size=10.5, bold=True, color=GREY)
    bullets(s, 6.55, 5.02, 6.2, 1.5, [
        "Do friends rate them much higher than strangers?",
        "Did friends rush to rate them right after a penalty?",
        "Has any stranger ever rated them at all?",
    ], size=11.5, gap=5, mark="·")

    # ══════════════════════════════════════════════ 4 farming — verified
    n += 1
    s = slide(prs, "Mechanism 1  ·  verified", "Volume Cannot Buy A Reputation",
              "Measured on the running system.", page=n + 1)
    figure(s, "fig_farming.png", 1.95, 3.70)
    text(s, M, 5.95, W, 0.9,
         "Ten fake five-star ratings from one friend moved a runner from 3.50 to "
         "3.55.  The attack was already worthless before any detection ran.",
         size=13, bold=True)
    text(s, M, 6.45, W, 0.5,
         "Penalty if flagged: the farmed reputation is discounted to nearly "
         "nothing, and the flag adds 700 m — fading to zero over 30 days.",
         size=11.5, color=GREY)

    # ═════════════════════════════════════════════════ 5 collusion rings
    n += 1
    s = slide(prs, "Mechanism 2", "Collusion Rings",
              "Farming done by a group: friends running errands for each other "
              "so the same money circles the same few people.", page=n + 1)
    figure(s, "fig_social.png", 2.05, 4.35, max_w=6.6, centre=False)
    text(s, 7.55, 2.00, 5.2, 0.3, "THE GAP", size=10.5, bold=True, color=GREY)
    text(s, 7.55, 2.30, 5.2, 1.5,
         "Viswanath et al., ACM SIGCOMM 2010 shows the social-graph defences are "
         "really community detection, and work only if the graph is a faithful "
         "observation of reality.\n\n"
         "In every case the graph is something the platform WATCHES.",
         size=11.5, spacing=1.2)
    text(s, 7.55, 4.20, 5.2, 0.3, "WHAT WE BUILT", size=10.5, bold=True, color=GREY)
    bullets(s, 7.55, 4.50, 5.2, 2.0, [
        "A ring is a closed loop of payments between mutual friends",
        "Three or more people. The loop went round more than once. Each leg is "
        "real — big enough, or often enough.",
        "Found with Tarjan's algorithm over the payment graph",
    ], size=11.5, gap=7)

    # ══════════════════════════ 6 how do you separate honest from insular
    n += 1
    s = slide(prs, "Mechanism 2  ·  the hard part",
              "Telling An Honest Group From A Closed One",
              "A friend group and a ring are the same shape. No single test "
              "separates them — each filter clears a different kind of innocent "
              "group.", page=n + 1)
    table(s, M, 2.15, W,
          ["", "What it asks", "Clears", "Still cannot separate"],
          [["1.  Structure", "Is the group closed?",
            "**nobody — a friend group and a ring both score 1.0",
            "everyone"],
           ["2.  Money cycle", "Does value ever leave the group?",
            "ordinary friend groups, who also run errands for strangers",
            "a group who genuinely only use each other"],
           ["3.  z", "Did our own app cause it?",
            "insular groups who simply accept what they are offered",
            "a group who deliberately prefer each other"],
           ["4.  Language model", "Do the errands read real?",
            "genuine roommates — different shops, odd hours, real notes",
            "a careful ring that writes varied titles"],
           ["5.  A human", "—", "—", "—"]],
          widths=[1.5, 2.9, 4.1, 3.63], hdr=10.5, body=10.5, rh=0.60)
    text(s, M, 5.85, W, 0.9,
         "Filter 1 is where naive detection stops — and it flags 6 of 6 entirely "
         "honest groups in simulation.", size=12.5, bold=True)
    text(s, M, 6.40, W, 0.6,
         "Whoever survives all four gets a person, not a ban.",
         size=12.5, italic=True, color=GREY)

    # ═══════════════════════════════════════════════════════ 7 the flaw
    n += 1
    s = slide(prs, "The problem we found", "Our Own Matcher Blinds Our Own Detector",
              "Cheng, Chen & Ye (IEEE ICDE 2019) co-assign people who have worked "
              "together before, and never examine that the policy decides who "
              "transacts with whom. Ours does the same.", page=n + 1)
    figure(s, "fig_boost.png", 2.15, 4.20)
    text(s, M, 6.48, W, 0.6,
         "The ring never changes its behaviour. Our routing catches up to it — "
         "and at our own setting there is nothing left to notice.",
         size=12.5, bold=True)

    # ════════════════════════════════════════════════════ 8 the offer log
    n += 1
    s = slide(prs, "Mechanism 3", "The Offer Log",
              "Maranzato & Pereira (IEEE LA-Web 2009) score sellers after the "
              "fact. Nothing records why the platform made the choice it made — "
              "so the platform's own contribution can never be subtracted.",
              page=n + 1)
    figure(s, "fig_offerlog.png", 2.15, 3.95)
    text(s, M, 6.30, W, 0.8,
         "Every errand, the app writes down who was in the running, why each "
         "ranked where they did, and who took it. Saving only the final score "
         "would make this question unanswerable.", size=12, color=GREY)

    # ═══════════════════════════════════════════════ 9 offer log verified
    n += 1
    s = slide(prs, "Mechanism 3  ·  verified", "What The Log Guarantees",
              page=n + 1)
    table(s, M, 2.00, W, ["Design decision", "Why", "Verified"],
          [["Save every part, not just the score",
            "A total cannot be taken apart. “Would he have won without the "
            "bonus?” is unanswerable from −55 alone.",
            "One shared formula for ranker and log, enforced by test — they "
            "cannot drift"],
           ["One record per offer round",
            "A re-offer happens to a different set of people. Merging rounds "
            "averages away the variation.",
            "153 rounds recorded on the running system"],
           ["Writing it can never break an errand",
            "Analytics sitting on a dispatch path must not stop somebody's "
            "dinner arriving.",
            "Test sabotages the writer — errand still created, still acceptable"]],
          widths=[3.0, 5.0, 4.13], hdr=10.5, body=10.5, rh=0.85)
    text(s, M, 5.55, W, 0.3, "THE BUG THIS CAUGHT", size=10.5, bold=True, color=GREY)
    text(s, M, 5.85, W, 1.0,
         "On the socially-blind rounds we ranked people without the friend "
         "bonus, but wrote down the score with it. The record contradicted what "
         "actually happened — and every later estimate would have inherited "
         "that, silently.\n"
         "Found by running it against the live stack. All 213 tests were "
         "passing.", size=12, spacing=1.25)

    # ══════════════════════════════════════════ 10 exploration — the text
    n += 1
    s = slide(prs, "Mechanism 4", "Deliberate Exploration",
              "1 in 20 errands, we ignore friendship completely.", page=n + 1)
    text(s, M, 2.00, 5.6, 0.3, "THE GAP", size=10.5, bold=True, color=GREY)
    bullets(s, M, 2.30, 5.6, 2.3, [
        "Perdomo et al., ICML 2020 | a deployed model changes the very "
        "distribution it is trying to predict.",
        "The standard remedy | collect some data under a policy you did not "
        "optimise, so you have something unbiased to compare against.",
        "Never applied here | exploration is used for unbiased evaluation in "
        "recommendation. Nobody has proposed it to keep fraud detectable.",
    ], size=11.5, gap=9)

    text(s, 6.65, 2.00, 6.1, 0.3, "WHAT WE BUILT", size=10.5, bold=True,
         color=GREY)
    bullets(s, 6.65, 2.30, 6.1, 2.3, [
        "5% of errands are ranked on distance alone | no friend bonus at all",
        "And no hop restriction | strangers see it immediately instead of "
        "waiting 45 seconds",
        "Both had to go together | removing the bonus while still hiding the "
        "errand from strangers leaves the sample exactly as biased",
    ], size=11.5, gap=9)

    text(s, M, 4.85, W, 0.3, "WHAT IT BUYS", size=10.5, bold=True, color=GREY)
    table(s, M, 5.15, 7.9,
          ["At our 1500 m setting", "The app hands them", "The ring takes",
           "Can we see it?"],
          [["no exploration", "99", "100", "**no — only +1"],
           ["**5% blind", "**96", "**100", "**yes — +4"],
           ["10% blind", "93", "100", "yes — +7"]],
          widths=[2.5, 1.9, 1.7, 1.8], hdr=10.5, body=10.5, rh=0.36)
    text(s, 8.75, 5.15, 4.0, 1.2,
         "The cost: 1 errand in 20 gets a worse-matched runner.\n\n"
         "It is a configuration value, not a redesign — it can be set to zero.",
         size=11.5, color=GREY, spacing=1.25)

    # ══════════════════════════════════ 11 exploration — where it sits
    n += 1
    s = slide(prs, "Mechanism 4", "Where The 5% Sits In Matching", page=n + 1)
    figure(s, "fig_matching.png", 1.55, 4.95)
    text(s, M, 6.62, W, 0.5,
         "The exploring branch is the control group: the only rounds whose "
         "outcome was not already shaped by friendship.",
         size=12, color=GREY)

    # ══════════════════════════════════════════════════════════ 11 the z
    n += 1
    s = slide(prs, "Mechanism 4  ·  the check", "What The App Expected, "
              "Versus What Happened", page=n + 1)
    text(s, M, 1.95, W, 0.4,
         "Out of 200 errands posted by members of the group:", size=12)
    table(s, M, 2.35, W,
          ["", "Rounds", "The app expected", "Honest group took", "Ring took"],
          [["Normal rounds", "190", "178.6", "177", "**190"],
           ["Socially blind rounds", "10", "3.0", "3", "**10"],
           ["**Total", "**200", "**181.6", "**180", "**200"]],
          widths=[3.2, 1.6, 2.9, 2.4, 2.03], hdr=11, body=11, rh=0.42)
    table(s, M, 4.05, W,
          ["", "Expected", "Observed", "Gap", "Verdict"],
          [["Honest group", "181.6", "180", "−1.6",
            "**our own routing explains it — no flag"],
           ["Ring", "181.6", "200", "+18.4",
            "**unexplained — continue to the ring check"]],
          widths=[2.6, 1.8, 1.8, 1.5, 4.43], hdr=11, body=11, rh=0.44)
    text(s, M, 5.85, W, 1.0,
         "Same expectation for both, because the app treats them identically. "
         "Only their behaviour differs — and it shows up most clearly in the "
         "blind rounds: 3 of 10 against 10 of 10.\n"
         "Ten blind errands carry about 40% of the evidence. That is why 5% is "
         "enough.", size=12, spacing=1.25)

    # ══════════════════════════════════════════════════════════ 12 the LLM
    n += 1
    s = slide(prs, "Mechanism 5", "Language Model Corroboration", page=n + 1)
    figure(s, "fig_llm.png", 1.70, 4.55)
    text(s, M, 6.42, W, 0.7,
         "Asked only about a group that already failed the money check and the z "
         "check. It reads what the errands say — the one question structure and "
         "money cannot answer.", size=12, color=GREY)

    # ═══════════════════════════════════════════════════ 13 penalisation
    n += 1
    s = slide(prs, "Consequences", "What A Flag Actually Does", page=n + 1)
    figure(s, "fig_penalty.png", 1.62, 4.70)
    text(s, M, 6.45, W, 0.6,
         "Severity is fixed per rule and never model-decided. If it responded to "
         "the model, an attacker could tune their errand titles to move their "
         "own penalty.", size=11.5, color=GREY)

    # ═══════════════════════════════════════════════════════ 14 summary
    n += 1
    s = slide(prs, "Summary", "Gaps Closed", page=n + 1)
    table(s, M, 2.00, W, ["Prior work", "The gap it leaves", "What we do"],
          [["Kamvar et al., WWW 2003",
            "Admits collusive groups defeat it; does not solve it",
            "Rings found as closed payment cycles among mutual friends"],
           ["Chiou & Tu, IEEE Access 2020",
            "Closeness only ever raises a rater's weight",
            "Closeness LOWERS it once ratings concentrate in one circle"],
           ["Viswanath et al., SIGCOMM 2010",
            "The graph is treated as observed, never as produced",
            "Ours is produced by our own routing — so evidence is conditioned "
            "on it"],
           ["Cheng, Chen & Ye, ICDE 2019",
            "The policy decides who transacts, and that is never examined",
            "That feedback loop is our central finding"],
           ["Maranzato & Pereira, LA-Web 2009",
            "Detection is retrospective and never reaches the matcher",
            "A flag is a bounded ranking penalty, applied from the next errand"],
           ["Perdomo et al., ICML 2020",
            "Performative prediction, applied to pricing — never to fraud",
            "The routing policy becomes the null hypothesis"]],
          widths=[2.9, 4.6, 4.63], hdr=10.5, body=10.5, rh=0.56)
    rule(s, 6.35)
    text(s, M, 6.55, W, 0.5,
         "234 automated tests passing   ·   ring detection, offer log and "
         "exploration verified on the running stack   ·   comparison figures "
         "from simulation", size=11.5, color=GREY)

    prs.save(OUT)
    print(f"written: {OUT}")
    print(f"  slides: {len(prs.slides._sldIdLst)}")


if __name__ == "__main__":
    build()
