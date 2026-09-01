"""Chris's Review-2 slides.

    cd docs/report && python deck_chris.py

Plain white, black text, no decoration. Every slide follows one shape, because
that is what the panel is marking:

    the gap in prior work  ->  what we built  ->  proof  ->  what it costs

Slide copy stays thin on purpose. Thresholds are answers to questions, not the
questions themselves — they live in the speaker's mouth, not on the screen. A
slide that contains everything you were going to say means the room reads
instead of listening.
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).parent
FIGS = HERE / "figures"
OUT = HERE / "Errandly-Review2-Chris.pptx"

INK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
RULE = RGBColor(0xCC, 0xCC, 0xCC)
HDR = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

M = 0.60           # left margin
W = 12.13          # usable width
TITLE_F = "Cambria"
BODY_F = "Calibri"
FOOTER = "Errandly  ·  Chris Martin Mattam  ·  23BCE0743"


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def text(slide, x, y, w, h, body, size=12, bold=False, color=INK,
         font=BODY_F, italic=False, align=PP_ALIGN.LEFT, spacing=1.15):
    tf = _tb(slide, x, y, w, h)
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tf


def bullets(slide, x, y, w, h, items, size=13, gap=8, mark="—"):
    tf = _tb(slide, x, y, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        lead, _, rest = item.partition("|")
        r = p.add_run()
        r.text = f"{mark}  "
        r.font.name = BODY_F
        r.font.size = Pt(size)
        r.font.color.rgb = GREY
        if rest:
            b = p.add_run()
            b.text = lead.strip() + "  "
            b.font.name = BODY_F
            b.font.size = Pt(size)
            b.font.bold = True
            b.font.color.rgb = INK
            t = p.add_run()
            t.text = rest.strip()
            t.font.name = BODY_F
            t.font.size = Pt(size)
            t.font.color.rgb = INK
        else:
            t = p.add_run()
            t.text = lead.strip()
            t.font.name = BODY_F
            t.font.size = Pt(size)
            t.font.color.rgb = INK
    return tf


def slide(prs, eyebrow, title, lede=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if eyebrow:
        text(s, M, 0.34, W, 0.24, eyebrow.upper(), size=10, bold=True, color=GREY)
    text(s, M, 0.58, W, 0.60, title, size=27, bold=True, font=TITLE_F)
    if lede:
        text(s, M, 1.36, W, 0.60, lede, size=12.5, color=GREY, spacing=1.2)
    text(s, M, 7.02, 9.0, 0.25, FOOTER, size=9, color=GREY)
    if page is not None:
        text(s, 11.6, 7.02, 1.1, 0.25, str(page), size=9, color=GREY,
             align=PP_ALIGN.RIGHT)
    return s


def table(slide, x, y, w, headers, rows, widths, hdr=9.5, body=9.5, rh=0.34):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(x), Inches(y), Inches(w),
                                   Inches(rh * (len(rows) + 1)))
    t = shape.table
    for i, cw in enumerate(widths):
        t.columns[i].width = Inches(cw)
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = HDR
        c.margin_left = c.margin_right = Inches(0.06)
        c.margin_top = c.margin_bottom = Inches(0.04)
        r = c.text_frame.paragraphs[0].add_run()
        r.text = h
        r.font.name = BODY_F
        r.font.size = Pt(hdr)
        r.font.bold = True
        r.font.color.rgb = INK
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE
            c.margin_left = c.margin_right = Inches(0.06)
            c.margin_top = c.margin_bottom = Inches(0.035)
            c.vertical_anchor = MSO_ANCHOR.TOP
            for li, line in enumerate(str(val).split("\n")):
                p = (c.text_frame.paragraphs[0] if li == 0
                     else c.text_frame.add_paragraph())
                p.line_spacing = 1.05
                r = p.add_run()
                r.text = line.replace("**", "")
                r.font.name = BODY_F
                r.font.size = Pt(body)
                r.font.bold = line.startswith("**")
                r.font.color.rgb = INK
    return t


def figure(slide, name, y, max_h, max_w=W, centre=True):
    """Place a figure scaled to fit a box, centred horizontally.

    Sized by HEIGHT as well as width. Picking a width by eye is how a figure
    ends up running off the bottom of the slide and under its own caption —
    which is exactly what happened on the first render.
    """
    from PIL import Image

    path = FIGS / name
    if not path.exists():
        raise FileNotFoundError(path)
    with Image.open(path) as im:
        aspect = im.width / im.height

    w = min(max_w, max_h * aspect)
    h = w / aspect
    x = (13.333 - w) / 2 if centre else M
    slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))
    return y + h


def rule(slide, y, x=M, w=W):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.012))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln
