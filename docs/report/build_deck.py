"""Build the Review-2 slide deck.

    cd docs/report && python build_deck.py

Reproduces the design system already used in Errandly-Review2-v2.pptx — 16:9,
Cambria titles over Calibri body, a grey eyebrow above each title and a footer
rule — so the revision is visually continuous with what was shown before.

Content follows the Review-1 submission: same title, same aim, same team. The
literature is cut to ten papers, each ending in the gap it leaves and what this
project does about it, and two slides are added that explain in plain terms what
a collusion ring is and what the project contributes.
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).parent
FIGS = HERE / "figures"
OUT = HERE / "Errandly-Review2-v3.pptx"

INK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
RULE = RGBColor(0xD8, 0xD8, 0xD8)
HDR_FILL = RGBColor(0xEF, 0xEF, 0xEF)
ACCENT_FILL = RGBColor(0xF7, 0xF7, 0xF7)

M = 0.62          # left margin
W = 12.09         # content width
TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

FOOTER = "Errandly  ·  School of Computer Science and Engineering"


def textbox(slide, x, y, w, h, text, size, bold=False, color=INK,
            font=BODY_FONT, align=PP_ALIGN.LEFT, italic=False, spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=12.5, color=INK, gap=6, bullet="—"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        head, _, tail = item.partition("|")
        r = p.add_run()
        r.text = f"{bullet}  "
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = GREY
        if tail:
            b = p.add_run()
            b.text = head.strip() + "   "
            b.font.name = BODY_FONT
            b.font.size = Pt(size)
            b.font.bold = True
            b.font.color.rgb = INK
            t = p.add_run()
            t.text = tail.strip()
            t.font.name = BODY_FONT
            t.font.size = Pt(size)
            t.font.color.rgb = color
        else:
            t = p.add_run()
            t.text = head.strip()
            t.font.name = BODY_FONT
            t.font.size = Pt(size)
            t.font.color.rgb = color
    return tb


def slide(prs, eyebrow, title, lede=None, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for shp in list(s.shapes):        # blank the layout placeholders
        shp._element.getparent().remove(shp._element)
    if eyebrow:
        textbox(s, M, 0.36, W, 0.24, eyebrow.upper(), 10, bold=True, color=GREY)
    textbox(s, M, 0.60, W, 0.62, title, 28, bold=True, color=INK, font=TITLE_FONT)
    y = 1.42
    if lede:
        textbox(s, M, y, W, 0.62, lede, 12, color=GREY, spacing=1.2)
    textbox(s, M, 7.00, 8.0, 0.25, FOOTER, 9, color=GREY)
    if page is not None:
        textbox(s, 11.71, 7.00, 1.0, 0.25, str(page), 9, color=GREY,
                align=PP_ALIGN.RIGHT)
    return s


def table(slide, x, y, w, headers, rows, widths, hdr_size=9.1, body_size=8.6,
          row_h=0.32):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(x), Inches(y), Inches(w),
                                   Inches(row_h * (len(rows) + 1)))
    t = shape.table
    for i, cw in enumerate(widths):
        t.columns[i].width = Inches(cw)
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = HDR_FILL
        c.margin_left = c.margin_right = c.margin_top = c.margin_bottom = Inches(0.056)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = BODY_FONT; r.font.size = Pt(hdr_size)
        r.font.bold = True; r.font.color.rgb = INK
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            c.margin_left = c.margin_right = Inches(0.056)
            c.margin_top = c.margin_bottom = Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.TOP
            for li, line in enumerate(str(val).split("\n")):
                p = c.text_frame.paragraphs[0] if li == 0 else c.text_frame.add_paragraph()
                p.line_spacing = 1.05
                r = p.add_run(); r.text = line
                r.font.name = BODY_FONT; r.font.size = Pt(body_size)
                r.font.color.rgb = INK
                if line.startswith("Ours:"):
                    r.font.bold = True
    return t


def figure(slide, name, x, y, w):
    path = FIGS / name
    if not path.exists():
        raise FileNotFoundError(path)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def panel(slide, x, y, w, h, heading, text, hsize=13, tsize=11.5):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT_FILL
    box.line.color.rgb = RULE; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = heading
    r.font.name = BODY_FONT; r.font.size = Pt(hsize); r.font.bold = True
    r.font.color.rgb = INK
    for line in text.split("\n"):
        p2 = tf.add_paragraph(); p2.line_spacing = 1.2; p2.space_before = Pt(4)
        r2 = p2.add_run(); r2.text = line
        r2.font.name = BODY_FONT; r2.font.size = Pt(tsize); r2.font.color.rgb = INK
    return box
